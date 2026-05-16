#!/usr/bin/env python3
"""
PPT + 音频 → 视频生成器
自动匹配播客转录与 PPT 页面，生成带字幕的视频
"""

import os
import sys
import json
import subprocess
import argparse
from pathlib import Path
from pptx import Presentation
import whisper
import yaml
from openai import OpenAI


class PPTVideoMaker:
    def __init__(self, config_path="config.yaml", backend="auto"):
        """初始化

        Args:
            config_path: 配置文件路径
            backend: AI 后端选择 ("auto", "claude", "codex", "custom")
        """
        self.backend = backend
        self.config = self._load_config(config_path)

        # 根据后端类型初始化客户端
        if self.backend == "codex":
            self.client = None  # Codex 使用 CLI，不需要客户端
            self.model = "codex"
        else:
            self.client = OpenAI(
                api_key=self.config["api"]["api_key"],
                base_url=self.config["api"]["base_url"]
            )
            self.model = self.config["api"]["model"]

    def _load_config(self, config_path):
        """加载配置文件，优先使用 Claude Code 的配置"""
        # 尝试从 Claude Code settings.json 读取配置
        claude_settings_path = os.path.expanduser("~/.claude/settings.json")
        api_config = None

        if os.path.exists(claude_settings_path):
            try:
                with open(claude_settings_path) as f:
                    settings = json.load(f)
                    env = settings.get("env", {})
                    if "ANTHROPIC_AUTH_TOKEN" in env and "ANTHROPIC_BASE_URL" in env:
                        api_config = {
                            "base_url": env["ANTHROPIC_BASE_URL"],
                            "api_key": env["ANTHROPIC_AUTH_TOKEN"],
                            "model": "claude-sonnet-4-20250514"
                        }
                        print("✅ 使用 Claude Code 配置")
            except Exception as e:
                print(f"⚠️  读取 Claude Code 配置失败: {e}")

        # 如果没有 Claude Code 配置，尝试读取本地配置文件
        if not api_config:
            if not os.path.exists(config_path):
                # 创建默认配置
                default_config = {
                    "api": {
                        "base_url": "https://api.openai.com/v1",
                        "api_key": "your-api-key-here",
                        "model": "gpt-4"
                    },
                    "whisper": {
                        "model": "medium",
                        "language": "zh"
                    },
                    "video": {
                        "resolution": "1920x1080",
                        "fps": 30
                    }
                }
                with open(config_path, "w") as f:
                    yaml.dump(default_config, f, allow_unicode=True)
                print(f"❌ 请先配置 {config_path} 文件")
                sys.exit(1)

            with open(config_path) as f:
                config = yaml.safe_load(f)
                api_config = config.get("api", {})

        # 合并配置
        if os.path.exists(config_path):
            with open(config_path) as f:
                local_config = yaml.safe_load(f)
        else:
            local_config = {
                "whisper": {
                    "model": "medium",
                    "language": "zh"
                },
                "video": {
                    "resolution": "1920x1080",
                    "fps": 30
                }
            }

        local_config["api"] = api_config
        return local_config

    def ppt_to_images(self, pptx_path, output_dir):
        """PPT 转图片（先转 PDF，再转图片）"""
        print(f"📄 正在转换 PPT 为图片...")
        output_dir = Path(output_dir)
        output_dir.mkdir(parents=True, exist_ok=True)

        # Step 1: PPT → PDF
        pdf_path = Path(pptx_path).with_suffix('.pdf')
        cmd = [
            "soffice",
            "--headless",
            "--convert-to", "pdf",
            str(pptx_path)
        ]

        result = subprocess.run(cmd, capture_output=True, text=True)
        if result.returncode != 0:
            print(f"❌ PPT 转 PDF 失败: {result.stderr}")
            sys.exit(1)

        # Step 2: PDF → PNG
        try:
            from pdf2image import convert_from_path
            images_pil = convert_from_path(str(pdf_path))

            image_paths = []
            for i, image in enumerate(images_pil):
                output_path = output_dir / f"slide_{i+1:03d}.png"
                image.save(output_path, "PNG")
                image_paths.append(output_path)

            print(f"✅ 已生成 {len(image_paths)} 张图片")

            # 清理临时 PDF
            pdf_path.unlink()

            return image_paths

        except ImportError:
            print(f"❌ 缺少 pdf2image 库，请运行: pip install pdf2image")
            sys.exit(1)

    def extract_ppt_text(self, pptx_path):
        """提取 PPT 每页的文字内容"""
        print(f"📝 正在提取 PPT 文字...")
        prs = Presentation(pptx_path)
        slides_text = []

        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text)

            slide_content = " | ".join(texts)
            slides_text.append({
                "slide_num": i + 1,
                "content": slide_content
            })
            print(f"  第 {i+1} 页: {slide_content[:50]}...")

        return slides_text

    def transcribe_audio(self, audio_path):
        """使用 Whisper 转录音频（带时间戳）"""
        print(f"🎤 正在转录音频（这可能需要几分钟）...")

        model = whisper.load_model(self.config["whisper"]["model"])
        result = model.transcribe(
            str(audio_path),
            language=self.config["whisper"]["language"],
            word_timestamps=True,
            verbose=False
        )

        # 提取带时间戳的句子
        segments = []
        for segment in result["segments"]:
            segments.append({
                "start": segment["start"],
                "end": segment["end"],
                "text": segment["text"]
            })

        print(f"✅ 转录完成，共 {len(segments)} 个片段")
        return segments, result

    def match_slides_with_transcript(self, slides_text, transcript_segments):
        """使用 AI 匹配 PPT 页面与转录时间点"""
        print(f"🤖 正在智能匹配 PPT 页面与音频...")

        # 构建 prompt
        slides_info = "\n".join([
            f"第 {s['slide_num']} 页: {s['content']}"
            for s in slides_text
        ])

        transcript_info = "\n".join([
            f"[{seg['start']:.1f}s - {seg['end']:.1f}s] {seg['text']}"
            for seg in transcript_segments[:50]  # 只发前 50 段避免太长
        ])

        prompt = f"""你是一个视频编辑助手。现在有一段播客音频的转录（带时间戳），以及对应的 PPT 每页内容。
请分析播客在哪个时间点开始讲下一页 PPT 的内容，输出 JSON 时间线。

## PPT 内容
{slides_info}

## 播客转录（前 50 段）
{transcript_info}

## 要求
1. 分析每页 PPT 的关键词与转录内容的对应关系
2. 确定每页 PPT 应该在哪个时间点开始显示
3. 输出 JSON 格式的时间线，格式如下：

```json
[
  {{"slide": 1, "start": 0.0, "end": 35.2}},
  {{"slide": 2, "start": 35.2, "end": 72.8}},
  ...
]
```

只输出 JSON，不要其他解释。"""

        response = self.client.chat.completions.create(
            model=self.model,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.3
        )

        # 解析 JSON
        content = response.choices[0].message.content
        # 提取 JSON（可能被包在 ```json ``` 中）
        if "```json" in content:
            content = content.split("```json")[1].split("```")[0]
        elif "```" in content:
            content = content.split("```")[1].split("```")[0]

        timeline = json.loads(content.strip())
        print(f"✅ 匹配完成，共 {len(timeline)} 个时间点")

        return timeline

    def generate_video(self, images, audio_path, timeline, output_path):
        """使用 FFmpeg 生成视频"""
        print(f"🎬 正在生成视频...")

        # 创建临时目录
        temp_dir = Path("temp_video_parts")
        temp_dir.mkdir(exist_ok=True)

        # 为每个 slide 生成视频片段
        video_parts = []
        for i, entry in enumerate(timeline):
            slide_num = entry["slide"]
            start = entry["start"]
            end = entry["end"]
            duration = end - start

            # 找到对应的图片
            image_path = images[slide_num - 1]
            part_path = temp_dir / f"part_{i:03d}.mp4"

            # 生成视频片段
            cmd = [
                "ffmpeg", "-y",
                "-loop", "1",
                "-i", str(image_path),
                "-t", str(duration),
                "-vf", f"scale={self.config['video']['resolution']}:force_original_aspect_ratio=decrease,pad={self.config['video']['resolution']}:(ow-iw)/2:(oh-ih)/2",
                "-r", str(self.config["video"]["fps"]),
                "-pix_fmt", "yuv420p",
                str(part_path)
            ]

            subprocess.run(cmd, capture_output=True, check=True)
            video_parts.append(part_path)
            print(f"  ✓ 第 {slide_num} 页 ({duration:.1f}s)")

        # 合并所有视频片段
        concat_file = temp_dir / "concat.txt"
        with open(concat_file, "w") as f:
            for part in video_parts:
                f.write(f"file '{part.absolute()}'\n")

        temp_video = temp_dir / "video_no_audio.mp4"
        cmd = [
            "ffmpeg", "-y",
            "-f", "concat",
            "-safe", "0",
            "-i", str(concat_file),
            "-c", "copy",
            str(temp_video)
        ]
        subprocess.run(cmd, capture_output=True, check=True)

        # 添加音频
        cmd = [
            "ffmpeg", "-y",
            "-i", str(temp_video),
            "-i", str(audio_path),
            "-c:v", "copy",
            "-c:a", "aac",
            "-shortest",
            str(output_path)
        ]
        subprocess.run(cmd, capture_output=True, check=True)

        print(f"✅ 视频已生成: {output_path}")

        # 清理临时文件
        for part in video_parts:
            part.unlink()
        concat_file.unlink()
        temp_video.unlink()
        temp_dir.rmdir()

    def add_subtitles(self, video_path, srt_path, output_path):
        """烧录字幕到视频"""
        print(f"📝 正在添加字幕...")

        cmd = [
            "ffmpeg", "-y",
            "-i", str(video_path),
            "-vf", f"subtitles={srt_path}",
            "-c:a", "copy",
            str(output_path)
        ]

        subprocess.run(cmd, capture_output=True, check=True)
        print(f"✅ 字幕已添加: {output_path}")

    def run(self, pptx_path, audio_path, output_path, add_subs=False):
        """完整流程"""
        print("=" * 60)
        print("PPT + 音频 → 视频生成器")
        print("=" * 60)

        # Step 1: PPT 转图片
        slides_dir = Path("slides")
        images = self.ppt_to_images(pptx_path, slides_dir)

        # Step 2: 提取 PPT 文字
        slides_text = self.extract_ppt_text(pptx_path)

        # Step 3: 转录音频
        transcript_segments, whisper_result = self.transcribe_audio(audio_path)

        # 保存 SRT 字幕文件
        srt_path = Path(audio_path).with_suffix(".srt")
        with open(srt_path, "w", encoding="utf-8") as f:
            for i, seg in enumerate(transcript_segments):
                f.write(f"{i+1}\n")
                f.write(f"{self._format_timestamp(seg['start'])} --> {self._format_timestamp(seg['end'])}\n")
                f.write(f"{seg['text']}\n\n")
        print(f"✅ 字幕文件已保存: {srt_path}")

        # Step 4: AI 匹配
        timeline = self.match_slides_with_transcript(slides_text, transcript_segments)

        # 保存时间线供人工检查
        timeline_path = Path("timeline.json")
        with open(timeline_path, "w", encoding="utf-8") as f:
            json.dump(timeline, f, indent=2, ensure_ascii=False)
        print(f"✅ 时间线已保存: {timeline_path}")
        print(f"💡 请检查 timeline.json，如需调整可手动修改后重新运行")

        # 询问是否继续
        response = input("\n是否继续生成视频？(y/n): ")
        if response.lower() != "y":
            print("已取消")
            return

        # Step 5: 生成视频
        self.generate_video(images, audio_path, timeline, output_path)

        # Step 6: 添加字幕（可选）
        if add_subs:
            final_output = Path(output_path).with_name(
                Path(output_path).stem + "_with_subs.mp4"
            )
            self.add_subtitles(output_path, srt_path, final_output)

        print("\n" + "=" * 60)
        print("✅ 全部完成！")
        print("=" * 60)

    def _format_timestamp(self, seconds):
        """格式化时间戳为 SRT 格式"""
        hours = int(seconds // 3600)
        minutes = int((seconds % 3600) // 60)
        secs = int(seconds % 60)
        millis = int((seconds % 1) * 1000)
        return f"{hours:02d}:{minutes:02d}:{secs:02d},{millis:03d}"


def main():
    parser = argparse.ArgumentParser(description="PPT + 音频 → 视频生成器")
    parser.add_argument("--pptx", required=True, help="PPT 文件路径")
    parser.add_argument("--audio", required=True, help="音频文件路径")
    parser.add_argument("--output", default="output.mp4", help="输出视频路径")
    parser.add_argument("--subs", action="store_true", help="是否添加字幕")
    parser.add_argument("--config", default="config.yaml", help="配置文件路径")

    args = parser.parse_args()

    maker = PPTVideoMaker(config_path=args.config)
    maker.run(args.pptx, args.audio, args.output, add_subs=args.subs)


if __name__ == "__main__":
    main()
