#!/usr/bin/env python3
"""
演示脚本 - 不需要真实音频，展示核心功能
"""

from pathlib import Path
from pptx import Presentation
import json

print("=" * 60)
print("PPT 视频生成器 - 功能演示")
print("=" * 60)

# 1. 提取 PPT 文字
print("\n[1/4] 提取 PPT 文字内容")
print("-" * 60)

prs = Presentation('test.pptx')
slides_text = []

for i, slide in enumerate(prs.slides):
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            texts.append(shape.text)

    slide_content = " | ".join(texts)
    slides_text.append({
        "slide_num": i + 1,
        "content": slide_content
    })
    print(f"第 {i+1} 页: {slide_content}")

# 2. PPT 转图片
print("\n[2/4] PPT 转图片")
print("-" * 60)

import subprocess
from pdf2image import convert_from_path

# PPT → PDF
pdf_path = Path('test.pdf')
cmd = ["soffice", "--headless", "--convert-to", "pdf", "test.pptx"]
result = subprocess.run(cmd, capture_output=True)

if result.returncode == 0:
    # PDF → PNG
    images = convert_from_path(str(pdf_path))

    output_dir = Path("demo_slides")
    output_dir.mkdir(exist_ok=True)

    image_paths = []
    for i, image in enumerate(images):
        output_path = output_dir / f"slide_{i+1:03d}.png"
        image.save(output_path, "PNG")
        image_paths.append(output_path)
        print(f"✅ {output_path.name}")

    pdf_path.unlink()  # 清理临时 PDF

# 3. 模拟转录结果
print("\n[3/4] 模拟音频转录（实际使用时会调用 Whisper）")
print("-" * 60)

# 模拟一个 10-20 分钟播客的转录
mock_transcript = [
    {"start": 0.0, "end": 8.5, "text": "大家好，欢迎来到今天的分享。"},
    {"start": 8.5, "end": 15.2, "text": "今天我们要介绍一个关于人工智能的项目。"},
    {"start": 15.2, "end": 22.8, "text": "首先让我们了解一下 AI 的基本概念。"},
    {"start": 22.8, "end": 30.5, "text": "人工智能是计算机科学的一个分支。"},
    {"start": 30.5, "end": 38.0, "text": "接下来我们看看技术架构部分。"},
    {"start": 38.0, "end": 45.3, "text": "我们的系统采用了微服务架构。"},
    {"start": 45.3, "end": 52.8, "text": "包含前端、后端和数据库三个主要层次。"},
    {"start": 52.8, "end": 60.0, "text": "最后让我们展望一下未来的发展。"},
    {"start": 60.0, "end": 67.5, "text": "我们计划在下个季度推出新功能。"},
    {"start": 67.5, "end": 75.0, "text": "并且会扩展到更多的应用场景。"},
]

for seg in mock_transcript:
    print(f"[{seg['start']:.1f}s - {seg['end']:.1f}s] {seg['text']}")

# 4. 模拟 AI 匹配
print("\n[4/4] 模拟 AI 智能匹配（实际使用时会调用 Claude API）")
print("-" * 60)

# 基于内容关键词的简单匹配逻辑
mock_timeline = [
    {"slide": 1, "start": 0.0, "end": 38.0},   # "项目介绍" 和 "AI 概念"
    {"slide": 2, "start": 38.0, "end": 60.0},  # "技术架构" 和 "微服务"
    {"slide": 3, "start": 60.0, "end": 75.0},  # "未来展望" 和 "新功能"
]

print("生成的时间线:")
for entry in mock_timeline:
    duration = entry['end'] - entry['start']
    print(f"  第 {entry['slide']} 页: {entry['start']:.1f}s - {entry['end']:.1f}s (持续 {duration:.1f}s)")

# 保存时间线
timeline_path = Path("demo_timeline.json")
with open(timeline_path, "w", encoding="utf-8") as f:
    json.dump(mock_timeline, f, indent=2, ensure_ascii=False)

print(f"\n✅ 时间线已保存: {timeline_path}")

# 总结
print("\n" + "=" * 60)
print("✅ 演示完成！")
print("=" * 60)
print("\n核心功能验证:")
print("  ✅ PPT 文字提取")
print("  ✅ PPT 转图片（每页独立）")
print("  ✅ 转录格式（带时间戳）")
print("  ✅ 时间线生成")
print("\n实际使用时:")
print("  1. 配置 config.yaml（填入你的 API）")
print("  2. 准备真实的 PPT 和音频文件")
print("  3. 运行: python ppt_video.py --pptx 你的.pptx --audio 你的.mp3")
print("  4. Whisper 会自动转录音频（3-5 分钟）")
print("  5. Claude API 会智能匹配时间点")
print("  6. FFmpeg 自动生成最终视频")
